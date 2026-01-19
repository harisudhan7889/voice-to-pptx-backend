import hmac
import json
import os
import time
from typing import List

import redis
from fastapi import FastAPI, Body, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from fastapi import Request
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
import hashlib
from pydantic import BaseModel
from starlette.responses import JSONResponse

FREE_LIMIT = 1

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"]
)


def get_redis_client() -> redis.Redis | None:
    url = os.getenv("REDIS_URL")
    if not url:
        print("REDIS_URL not set, guest limiting disabled")
        return None
    try:
        client = redis.from_url(url)
        # Light ping to validate connection
        client.ping()
        print("Connected to Redis successfully")
        return client
    except Exception as e:
        print(f"Redis connection failed: {e}")
        return None


# Redis connection
r = get_redis_client()


def create_guest_id(request: Request) -> str:
    """Fingerprint = IP + User-Agent (privacy-safe)"""
    ip = request.client.host.replace("::ffff:", "")
    ua = request.headers.get("user-agent", "")[:100]
    fingerprint = f"{ip}:{ua}"
    return hashlib.md5(fingerprint.encode()).hexdigest()


# @app.middleware("http")
# async def guest_limiter(request: Request, call_next):
#     # If Redis not available, skip limiting (but app still works)
#     if r is None:
#         return await call_next(request)
#
#     if request.url.path == "/api/generate-pptx" and request.method == "POST":
#         guest_id = create_guest_id(request)
#         key = f"guest:{guest_id}"
#
#         # Track usage
#         count = r.incr(key)
#         r.expire(key, 2592000)  # 30 days
#
#         # Store in request state for PPT generation
#         request.state.guest_count = count
#         request.state.guest_id = guest_id
#
#         if count > FREE_LIMIT:
#             return JSONResponse(
#                 status_code=402,
#                 content={
#                     "status": "limit_reached",
#                     "used": count,
#                     "limit": FREE_LIMIT,
#                     "guest_id": guest_id,
#                     "upgrade_url": "/upgrade",
#                     "message": "You've created 3 amazing presentations!"
#                 }
#             )
#
#     response = await call_next(request)
#     return response


@app.get("/api/templates")
async def get_templates():
    """
    Fetch available templates grouped by type
    """
    return {
        "status": "success",
        "templates": [
            {
                "id": "geometric",
                "name": "geometric",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_geometric_1.png", "/thumbnails/thumbnail_geometric_2.png"]
            },
            {
                "id": "streamline",
                "name": "Streamline",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_streamline_1.png", "/thumbnails/thumbnail_streamline_2.png"]
            },
            {
                "id": "swiss",
                "name": "swiss",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_swiss_1.png", "/thumbnails/thumbnail_swiss_2.png"]
            },
            {
                "id": "momentum",
                "name": "momentum",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_momentum_1.png", "/thumbnails/thumbnail_momentum_2.png"]
            },
            {
                "id": "material",
                "name": "material",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_material_1.png", "/thumbnails/thumbnail_material_2.png"]
            },
            {
                "id": "slate",
                "name": "slate",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_slate_1.png", "/thumbnails/thumbnail_slate_2.png"]
            },
            {
                "id": "paperback",
                "name": "paperback",
                "description": "",
                "thumbnails": ["/thumbnails/thumbnail_paperback_1.png", "/thumbnails/thumbnail_paperback_2.png"]
            }
        ]
    }


class Slide(BaseModel):
    slideNumber: int
    title: str
    type: str
    format: str
    content: List[str]


class SlidesRequest(BaseModel):
    templateId: str
    title: str
    slides: List[Slide]


@app.post("/api/generate-pptx")
async def generate_pptx(request: SlidesRequest = Body(...), request_obj: Request = None):
    try:
        # NEW: Check Pro status FIRST (middleware sets this)
        is_pro = getattr(request_obj.state, 'is_pro', False) if request_obj else False

        # Get guest info from middleware (if exists and NOT pro)
        guest_count = getattr(request_obj.state, 'guest_count', 0) if request_obj and not is_pro else 0
        is_guest = not is_pro and (not request_obj or not request_obj.headers.get("authorization"))

        slides_data = request.dict()
        template_id = slides_data.get("templateId")

        print(f" - Pro: {is_pro}, Guest count: {guest_count}")

        # Pro users = NO watermark, NO limits
        if is_pro:
            print("Pro user - Unlimited access")
        # Free users = watermark on #3+
        elif is_guest and guest_count >= FREE_LIMIT:
            print("Free user - Watermark added")

        template_files = {
            "geometric": "templates/geometric.pptx",
            "streamline": "templates/streamline.pptx",
            "swiss": "templates/swiss.pptx",
            "momentum": "templates/momentum.pptx",
            "material": "templates/material.pptx",
            "slate": "templates/slate.pptx",
            "paperback": "templates/paperback.pptx"
        }
        template_path = template_files.get(template_id, template_files["slate"])

        prs = Presentation(template_path)
        title_template = prs.slides[0]
        content_template = prs.slides[1]

        for slide_data in slides_data["slides"]:
            slide_type = slide_data.get("type", "content")
            template = title_template if slide_type == "title" else content_template
            slide = prs.slides.add_slide(template.slide_layout)
            add_slide_content(slide, slide_data)

        r_id_list = prs.slides._sldIdLst
        del r_id_list[1]
        del r_id_list[0]

        rc_user_id = request_obj.headers.get("X-RC-App-User-ID") if request_obj else None
        user_id = rc_user_id or getattr(request_obj.state, 'guest_id', f"guest-{int(time.time())}")

        # Watermark ONLY for free users (PPT #3+), NEVER for Pro
        if is_guest and guest_count >= FREE_LIMIT:
            add_watermark(prs, "Made with Voice-to-PPT Free")
        timestamp = int(time.time())
        filename = f"presentation-{user_id[:8]}-{template_id}-{timestamp}.pptx"
        # filename = f"presentation-{template_id}-{int(time.time())}.pptx"
        os.makedirs("presentations", exist_ok=True)
        output_path = f"presentations/{filename}"
        prs.save(output_path)

        ppt_info = {
            "filename": filename,
            "template": template_id,
            "created": timestamp,
            "url": f"/download/{filename}",
            "is_pro": is_pro
        }

        # Save to user's history list (Pro: permanent, Guest: session)
        history_key = f"history:{user_id}"
        r.lpush(history_key, json.dumps(ppt_info))

        # Pro: Keep 50 PPTs | Free: Keep 3 PPTs
        max_history = 50 if is_pro else 3
        r.ltrim(history_key, 0, max_history - 1)

        # Keep history alive (Pro: 30 days, Guest: 1 hour)
        ttl = 2592000 if is_pro else 3600  # 30 days vs 1 hour
        r.expire(history_key, ttl)

        return {
            "status": "success",
            "downloadUrl": f"/download/{filename}",
            "is_pro": is_pro,
            "is_guest": is_guest,
            "guest_count": guest_count,
            "user_id": user_id,
            "show_upgrade": is_guest and guest_count >= FREE_LIMIT,
            "watermark": is_guest and guest_count >= FREE_LIMIT
        }

    except Exception as e:
        return {"status": "error", "message": str(e)}


@app.get("/api/ppt-history")
async def get_ppt_history(request: Request):
    try:
        # Skip if no Redis
        if r is None:
            print("Redis unavailable - empty history")
            return {"ppt_history": [], "count": 0}

        # Get user ID safely (Pro OR Guest)
        rc_user_id = request.headers.get("X-RC-App-User-ID")
        if rc_user_id:
            history_key = f"history:{rc_user_id}"
        else:
            guest_id = create_guest_id(request)
            if not guest_id:
                print("No user ID available")
                return {"ppt_history": [], "count": 0}
            history_key = f"history:{guest_id}"

        # Get history from Redis (safe decode)
        history_raw = r.lrange(history_key, 0, 9)  # Max 10 recent

        # Parse JSON safely
        ppt_list = []
        for item in history_raw:
            try:
                ppt_info = json.loads(item.decode('utf-8', errors='ignore'))
                ppt_list.append(ppt_info)
            except (json.JSONDecodeError, UnicodeDecodeError) as e:
                print(f"Corrupt history item: {e}")
                continue  # Skip bad items

        print(f"History loaded: {len(ppt_list)} PPTs for {history_key[:20]}...")
        return {
            "ppt_history": ppt_list,
            "count": len(ppt_list)
        }

    except redis.RedisError as e:
        print(f"Redis error in history: {e}")
        return {"ppt_history": [], "count": 0}

    except Exception as e:
        print(f"Unexpected error in get_ppt_history: {e}")
        raise HTTPException(status_code=500, detail="History service temporarily unavailable")


def add_watermark(prs, text: str):
    """Add subtle watermark to all slides for free users"""
    for slide in prs.slides:
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(7)
        height = Inches(0.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xF5, 0xA6, 0x23)  # Gold color
        p.alignment = PP_ALIGN.CENTER


def get_correct_layout(prs, slide_type):
    """Get correct layout for slide type"""

    # Title slide - always first
    if slide_type == "title":
        return prs.slide_layouts[0]  # Title Slide

    # Content slide - look for layouts with title + content placeholders
    for layout in prs.slide_layouts[1:]:  # Skip title
        # Check for title (placeholder 0) + content (placeholder 1)
        has_title = False
        has_content = False

        for shape in layout.placeholders:
            if shape.placeholder_format.idx == 0:  # Title placeholder
                has_title = True
            elif shape.placeholder_format.idx == 1:  # Content placeholder
                has_content = True

        if has_title and has_content:
            return layout

    # Fallback: first layout with content
    for layout in prs.slide_layouts:
        if any(shape.placeholder_format.idx == 1 for shape in layout.placeholders):
            return layout

    return prs.slide_layouts[1]  # Final fallback


def add_slide_content(slide, slide_data):
    """Update text in placeholders - NEVER clear text_frame"""

    slide_type = slide_data.get("type", "content")
    title = slide_data.get("title", "")
    content_format = slide_data.get("format", "bullets")  # NEW: Get format type
    content = slide_data.get("content", [])

    print(f" - Content: {content_format}")

    # Get ALL text shapes in order
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    text_shapes.sort(key=lambda s: s.top)

    if slide_type == "title":
        # Title slide
        if len(text_shapes) >= 1:
            # Update first shape (title)
            p = text_shapes[0].text_frame.paragraphs[0]
            p.text = title

        if len(text_shapes) >= 2 and content:
            # Update second shape (subtitle)
            p = text_shapes[1].text_frame.paragraphs[0]
            if content_format == "paragraph":
                p.text = content[0]  # Single string for paragraph
            else:
                p.text = content[0]  # First item for bullets

    else:  # content slide
        # Content slide
        if len(text_shapes) >= 1:
            # Update first shape (title)
            p = text_shapes[0].text_frame.paragraphs[0]
            p.text = title

        if len(text_shapes) >= 2:
            # Update second shape (content)
            content_frame = text_shapes[1].text_frame

            # IMPORTANT: Don't clear! Just update/add paragraphs
            # Remove ALL paragraphs except first
            while len(content_frame.paragraphs) > 1:
                # Get the paragraph element
                p_element = content_frame.paragraphs[1].element
                # Remove from XML
                p_element.getparent().remove(p_element)

            # NEW: Handle based on format type
            if content_format == "paragraph":
                # Single paragraph content
                p = content_frame.paragraphs[0]
                p.text = content[0]  # content is a string
                p.level = 0
            else:  # bullets format
                # Multiple bullet points (content is array)
                for i, bullet in enumerate(content):
                    if i == 0:
                        # Update first paragraph (already exists)
                        p = content_frame.paragraphs[0]
                        p.text = bullet
                    else:
                        # Add new paragraph
                        p = content_frame.add_paragraph()
                        p.text = bullet

                    p.level = 0


def debug_slide(slide):
    """Print all shapes in slide"""
    for i, shape in enumerate(slide.shapes):
        print(f"Shape {i}: {shape.name}")
        if shape.has_text_frame:
            print(f" - Has text: {shape.text_frame.text[:50]}")
        else:
            print(f" - No text (image/design)")


@app.get("/debug")
async def debug():
    import os
    return {
        "cwd": os.getcwd(),
        "thumbnails_exists": os.path.exists("thumbnails"),
        "thumbnails_contents": os.listdir("thumbnails") if os.path.exists("thumbnails") else "NO FOLDER",
        "mount_status": "Check if app.mount line exists above"
    }


@app.get("/download/{filename}")
async def download_presentation(filename: str):
    return FileResponse(
        path=f"presentations/{filename}",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


async def verify_revenuecat_signature(auth_header: str = None) -> bool:
    """Verify RevenueCat webhook using Authorization header"""
    if not auth_header:
        return True  # Skip if no auth configured

    expected_token = os.getenv("REVENUECAT_WEBHOOK_SECRET")
    if not expected_token:
        return True

    # RevenueCat sends: Authorization: Bearer <your-secret>
    try:
        token = auth_header.replace("Bearer ", "")
        return hmac.compare_digest(token, expected_token)
    except:
        return False


@app.post("/api/revenuecat-webhook")
async def revenuecat_webhook(request: Request, payload: dict = Body(...)):
    # FIXED: Use Authorization header (RevenueCat standard)
    auth_header = request.headers.get("Authorization", "")
    if not await verify_revenuecat_signature(auth_header):
        print("Invalid webhook auth")
        raise HTTPException(status_code=401, detail="Invalid auth")

    print(f"RevenueCat webhook: {payload.get('event', {}).get('type', 'unknown')}")

    if r is None:
        print("Redis unavailable")
        return {"status": "redis_unavailable"}

    event = payload.get("event", {})
    if not event:
        return {"status": "invalid_payload"}

    event_type = event.get("type")
    app_user_id = event.get("app_user_id")
    product_id = event.get("product_id", "")

    if not all([event_type, app_user_id]):
        print(f"Missing fields: {event_type}, {app_user_id}")
        return {"status": "missing_fields"}

    pro_key = f"pro:{app_user_id}"

    try:
        if event_type == "TEST":
            print(f"🧪 TEST SUCCESS: {app_user_id}")
            return {"status": "test_success"}

        elif event_type == "INITIAL_PURCHASE":
            if not await r.exists(pro_key):
                if "lifetime" in product_id.lower() or "149" in product_id:
                    await r.set(pro_key, "lifetime")
                    print(f"LIFETIME: {app_user_id}")
                else:
                    await r.setex(pro_key, 31536000, "subscription")
                    print(f"SUBSCRIPTION: {app_user_id}")
            else:
                print(f"Already pro: {app_user_id}")

        elif event_type == "RENEWAL":
            await r.setex(pro_key, 31536000, "subscription")
            print(f"RENEWED: {app_user_id}")

        elif event_type == "CANCELLATION":
            # FIXED: Simple atomic check
            status = await r.get(pro_key)
            if status == b"subscription":
                await r.delete(pro_key)
                print(f"CANCELLED: {app_user_id}")

        elif event_type == "EXPIRATION":
            await r.delete(pro_key)
            print(f"EXPIRED: {app_user_id}")

        else:
            print(f"Unknown: {event_type}")

    except Exception as e:
        print(f"Webhook error: {e}")
        return {"status": "error"}

    return {"status": "ok"}


@app.middleware("http")
async def pro_guest_limiter(request: Request, call_next):
    if request.url.path != "/api/generate-pptx" or request.method != "POST":
        return await call_next(request)

    if r is None:
        print("Redis unavailable - skipping limits")
        return await call_next(request)

    try:
        # 1. ASYNC PRO CHECK (FIXED)
        rc_user_id = request.headers.get("X-RC-App-User-ID")
        if rc_user_id:
            pro_status = r.get(f"pro:{rc_user_id}")
            if pro_status:
                request.state.is_pro = True
                print(f"Pro: {rc_user_id} ({pro_status.decode(errors='ignore')})")
                return await call_next(request)

        # 2. ASYNC GUEST TRACKING (FIXED)
        guest_id = create_guest_id(request)
        count = r.incr(f"guest:{guest_id}")
        r.expire(f"guest:{guest_id}", 2592000)

        request.state.guest_count = int(count)
        request.state.guest_id = guest_id
        print(f"Guest #{count}: {guest_id[:8]}...")

        # 3. FIXED LIMIT RESPONSE
        if count > FREE_LIMIT:
            return JSONResponse(
                status_code=402,
                content={
                    "status": "limit_reached",
                    "used": int(count),
                    "limit": FREE_LIMIT,
                    "guest_id": guest_id,
                    "upgrade_url": "/upgrade",
                    "message": f"You've created {FREE_LIMIT} amazing presentations!"
                }
            )

    except Exception as e:
        print(f"Middleware error: {e}")

    return await call_next(request)


app.mount("/thumbnails", StaticFiles(directory="thumbnails"), name="thumbnails")

if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)  # ← host="0.0.0.0"
